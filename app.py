import streamlit as st
import os
import re
from dotenv import load_dotenv
from tavily import TavilyClient
from google import genai
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email.mime.text import MIMEText
from email import encoders
import tempfile
import logging

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

# Mock Data
MOCK_EMAIL = "Subject: CRM Issues. From: John @ [Prospect]. Hi, we are struggling with data silos. Our current tool is too slow."
MOCK_TRANSCRIPT = "Teams Meeting: We need AI features. Budget is around $50k/year. Need implementation in Q1."

# Page configuration
st.set_page_config(
    page_title="NexusCRM Sales Proposal Copilot",
    page_icon="🤖",
    layout="wide"
)

# Custom CSS for styling
st.markdown("""
<style>
    /* Remove default Streamlit padding */
    .block-container {
        padding-top: 1rem;
        padding-bottom: 0rem;
    }
    
    /* M365 background color */
    .stApp {
        background-color: #f3f2f1;
    }
    
    /* Font family */
    * {
        font-family: 'Segoe UI', -apple-system, BlinkMacSystemFont, sans-serif;
    }
    
    /* Hide Streamlit header and footer */
    header {visibility: hidden;}
    footer {visibility: hidden;}
    
    /* Chat message styling */
    .stChatMessage {
        background-color: white;
        border-radius: 8px;
        padding: 12px 16px;
        margin-bottom: 12px;
        box-shadow: 0 1px 2px rgba(0,0,0,0.08);
    }
    
    /* Input styling */
    .stChatInputContainer {
        border-top: 1px solid #e1dfdd;
        padding-top: 16px;
    }
    
    /* Right panel styling */
    .right-panel {
        background-color: #f5f5f5;
        padding: 20px;
        border-radius: 8px;
        border: 1px solid #e1dfdd;
        height: 100%;
    }
    
    /* Container styling for editable forms */
    .draft-container {
        background-color: white;
        border: 1px solid #e1dfdd;
        border-radius: 8px;
        padding: 16px;
        margin: 12px 0;
    }
</style>
""", unsafe_allow_html=True)

# Initialize Session State
if "messages" not in st.session_state:
    st.session_state.messages = []

if "company_data" not in st.session_state:
    st.session_state.company_data = {
        "name": "",
        "full_draft": "",
        "edited_full_draft": "",
        "ppt_theme": {
            "bg_color": [243, 242, 241],
            "title_color": [0, 120, 212],
            "body_color": [50, 49, 48],
            "accent_color": [0, 120, 212]
        }
    }

if "uploaded_files" not in st.session_state:
    st.session_state.uploaded_files = ["Company_Overview.pdf", "Pricing_Tier_2025.pptx"]

if "view_mode" not in st.session_state:
    st.session_state.view_mode = "main"  # or "email_view" or "chat_view"

if "company_emails" not in st.session_state:
    st.session_state.company_emails = []

if "company_chats" not in st.session_state:
    st.session_state.company_chats = []

if "show_send_modal" not in st.session_state:
    st.session_state.show_send_modal = False

if "send_modal_message_id" not in st.session_state:
    st.session_state.send_modal_message_id = None

if "last_sent_recipients" not in st.session_state:
    st.session_state.last_sent_recipients = []

if "last_sent_message_id" not in st.session_state:
    st.session_state.last_sent_message_id = None

if "selected_recipients" not in st.session_state:
    st.session_state.selected_recipients = []

# Helper Functions
def research_company(name, emails=None, chats=None):
    """Research company using Tavily and generate proposal content using the new Gemini SDK."""
    try:
        tavily = TavilyClient(api_key=os.getenv("TAVILY_API_KEY"))
        client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))

        # 1. Search Tavily
        logger.info(f"Searching Tavily for {name}...")
        search_query = f"{name} strategic goals 2025 financial challenges recent news"
        search_result = tavily.search(query=search_query, search_depth="advanced")
        context = search_result.get("results", [])
        
        # 2. Format emails for context
        email_context = ""
        if emails:
            email_summaries = []
            for e in emails[:5]:  # Use up to 5 emails
                email_summaries.append(f"From {e['sender']} ({e['date']}): {e['subject']}\n{e['body'][:300]}...")
            email_context = "\n\n".join(email_summaries)
        else:
            email_context = MOCK_EMAIL
        
        # 3. Format Teams chat for context
        chat_context = ""
        if chats and len(chats) > 0:
            chat = chats[0]
            chat_messages = []
            for m in chat['messages']:
                chat_messages.append(f"{m['sender']} ({m['timestamp']}): {m['content']}")
            chat_context = "\n".join(chat_messages)
        else:
            chat_context = MOCK_TRANSCRIPT

        # 4. Generate Content with Gemini
        logger.info(f"Generating proposal for {name} using Gemini (gemini-2.0-flash)...")
        prompt = f"""
        You are writing a sales proposal from 'NexusCRM' (a CRM software company) to {name}.
        
        Context gathered:
        - Web Research: {json.dumps(context)}
        - Recent Email Communications:
{email_context}
        
        - Teams Discussion:
{chat_context}
        
        Based on this information, draft a sales proposal with 3 distinct sections.
        The values for each key MUST be a plain string (markdown formatted), NOT a nested JSON object.
        
        1. **Executive Summary**: Brief overview addressing their pain points (data silos, speed issues).
        2. **The NexusCRM Solution**: How our AI-powered CRM can help with their specific needs.
        3. **Investment**: Pricing proposal aligned with their ~$50k/year budget, including implementation timeline for Q1.

        Return the output strictly as a JSON object with keys: "executive_summary", "solution", "pricing".
        """

        response = client.models.generate_content(
            model='gemini-2.0-flash',
            contents=prompt
        )
        content = response.text
        logger.info(f"Raw AI Response: {content}")
        
        if content.startswith("```json"):
            content = content.replace("```json", "").replace("```", "")
        if content.startswith("```"):
            content = content.replace("```", "")
            
        data = json.loads(content)
        
        # Ensure all values are strings (prevent nested JSON appearing in UI)
        for key in ["executive_summary", "solution", "pricing"]:
            if isinstance(data.get(key), (dict, list)):
                data[key] = json.dumps(data[key], indent=2)
                
        logger.info("Successfully parsed and cleaned response.")
        return data

    except Exception as e:
        logger.error(f"Error in research_company: {str(e)}", exc_info=True)
        st.error(f"NexusCRM Agent Error: {e}")
        return {
            "executive_summary": f"Research/AI Error: {str(e)}",
            "solution": "Check your .env for Tavily/Gemini API keys.",
            "pricing": "Internal Error."
        }

def extract_contacts(emails, chats):
    """Extract unique contacts from emails and Teams chats."""
    contacts = []
    
    # Extract from emails
    for email in emails:
        sender = email['sender']
        # Parse email to get name and address
        if '@' in sender:
            name_part = sender.split('@')[0].replace('.', ' ').title()
            contacts.append({
                'email': sender,
                'name': name_part,
                'source': 'Email'
            })
    
    # Extract from Teams chat participants
    if chats:
        for chat in chats:
            for participant in chat.get('participants', []):
                # Extract name from "Name (Role)" format
                name = participant.split('(')[0].strip()
                # Generate email from name
                email_name = name.lower().replace(' ', '.')
                company_domain = st.session_state.company_data.get('name', 'company').lower().replace(' ', '').replace(',', '')
                contacts.append({
                    'email': f"{email_name}@{company_domain}.com",
                    'name': name,
                    'source': 'Teams Chat'
                })
    
    # Remove duplicates based on email
    seen = set()
    unique_contacts = []
    for contact in contacts:
        if contact['email'] not in seen:
            seen.add(contact['email'])
            unique_contacts.append(contact)
    
    return unique_contacts

def render_send_modal(message_id):
    """Render the recipient selection modal for sending proposals."""
    if not st.session_state.show_send_modal:
        return
    if st.session_state.send_modal_message_id != message_id:
        return

    st.markdown("---")
    st.markdown("### 📧 Select Recipients")
    st.caption("These contacts were found in emails and Teams chats")

    contacts = extract_contacts(
        st.session_state.company_emails,
        st.session_state.company_chats
    )

    for i, contact in enumerate(contacts):
        is_selected = contact['email'] in st.session_state.selected_recipients
        
        col_check, col_info = st.columns([0.1, 0.9])
        with col_check:
            checked = st.checkbox(
                "", 
                value=is_selected, 
                key=f"checkbox_{contact['email']}_{message_id}"
            )
            if checked and contact['email'] not in st.session_state.selected_recipients:
                st.session_state.selected_recipients.append(contact['email'])
            elif not checked and contact['email'] in st.session_state.selected_recipients:
                st.session_state.selected_recipients.remove(contact['email'])
        
        with col_info:
            st.markdown(f"**{contact['name']}** - {contact['email']}")
            st.caption(f"Found in: {contact['source']}")

    st.markdown("---")

    btn_col1, btn_col2, btn_col3 = st.columns([1, 1, 1])
    
    with btn_col1:
        if st.button("Cancel", key=f"cancel_send_{message_id}"):
            st.session_state.show_send_modal = False
            st.session_state.send_modal_message_id = None
            st.rerun()
    
    with btn_col3:
        selected_count = len(st.session_state.selected_recipients)
        if st.button(
            f"📤 Send to {selected_count} Selected", 
            key=f"confirm_send_{message_id}", 
            type="primary",
            use_container_width=True,
            disabled=(selected_count == 0)
        ):
            if st.session_state.selected_recipients:
                recipients = list(st.session_state.selected_recipients)
                st.session_state.last_sent_recipients = recipients
                st.session_state.last_sent_message_id = message_id
                st.session_state.show_send_modal = False
                st.session_state.send_modal_message_id = None
                st.rerun()

def get_theme_update(user_suggestion, current_theme):
    """Use Gemini to translate a theme suggestion into RGB values."""
    client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))
    prompt = f"""
    Current PPT Theme (RGB):
    {json.dumps(current_theme)}
    
    User Suggestion: "{user_suggestion}"
    
    Translate this suggestion into a new RGB theme. 
    Return strictly a JSON object with these keys:
    - bg_color: [R, G, B]
    - title_color: [R, G, B]
    - body_color: [R, G, B]
    - accent_color: [R, G, B]
    """
    try:
        response = client.models.generate_content(
            model='gemini-2.0-flash',
            contents=prompt
        )
        content = response.text
        if "```json" in content:
            content = content.split("```json")[1].split("```")[0]
        elif "```" in content:
            content = content.split("```")[1].split("```")[0]
        return json.loads(content.strip())
    except:
        return current_theme

def generate_pptx(data):
    """Generate PowerPoint presentation based on content and theme."""
    prs = Presentation()
    theme = data.get('ppt_theme', {})
    
    def rgb_to_color(rgb_list):
        return RGBColor(rgb_list[0], rgb_list[1], rgb_list[2])

    bg_color = rgb_to_color(theme.get('bg_color', [255, 255, 255]))
    title_color = rgb_to_color(theme.get('title_color', [0, 120, 212]))
    body_color = rgb_to_color(theme.get('body_color', [0, 0, 0]))

    def add_slide(title_text, content_text):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        title.text_frame.paragraphs[0].font.bold = True
        
        content = slide.placeholders[1]
        content.text = content_text
        # Optional: Set body color if needed (more complex to iterate all runs)

    # Parsing logic
    draft = data.get('edited_full_draft', '')
    sections = {}
    current_key = "Executive Summary"
    for line in draft.split('\n'):
        if re.match(r'^(#|##|\*\*)\s*(Executive Summary|Understanding)', line, re.I):
            current_key = "Executive Summary"
            sections[current_key] = ""
        elif re.match(r'^(#|##|\*\*)\s*(Solution|The Nexus)', line, re.I):
            current_key = "Solution"
            sections[current_key] = ""
        elif re.match(r'^(#|##|\*\*)\s*(Investment|Pricing)', line, re.I):
            current_key = "Investment"
            sections[current_key] = ""
        else:
            if current_key:
                sections[current_key] = sections.get(current_key, "") + line + "\n"

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    title = slide.shapes.title
    title.text = f"NexusCRM → {data['name']}"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Strategic Proposal for Digital Transformation"

    # slides 2-4
    add_slide("Understanding Your Needs", sections.get("Executive Summary", "Details in full draft."))
    add_slide("The NexusCRM Solution", sections.get("Solution", "Details in full draft."))
    add_slide("Investment/Pricing", sections.get("Investment", "Details in full draft."))

    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(temp_file.name)
    return temp_file.name

# Main Layout: Two Columns
chat_col, right_panel = st.columns([0.7, 0.3])

# ===== RIGHT PANEL =====
with right_panel:
    st.markdown('<div class="right-panel">', unsafe_allow_html=True)
    
    # Tools Section - Show discovery results if company is searched
    st.markdown("### Context & Tools")
    
    if st.session_state.company_data.get("name"):
        # Show discovery results
        company_name = st.session_state.company_data["name"]
        
        st.success(f"📧 {len(st.session_state.company_emails)} emails found")
        if st.button("View Emails", key="view_emails_btn", use_container_width=True):
            st.session_state.view_mode = "email_view"
            st.rerun()
        
        st.success(f"💬 {len(st.session_state.company_chats)} group chat found")
        if st.button("View Chat", key="view_chat_btn", use_container_width=True):
            st.session_state.view_mode = "chat_view"
            st.rerun()
    else:
        # Default state before search
        st.info("✅ Outlook - Active")
        st.info("✅ Teams - Active")
    
    st.markdown("---")
    
    # Knowledge Section
    st.markdown("### Knowledge")
    uploaded_file = st.file_uploader(
        "Upload Past Proposals/Brochures",
        type=["pdf", "pptx"],
        help="Upload reference documents to improve proposal quality"
    )
    
    if uploaded_file:
        if uploaded_file.name not in st.session_state.uploaded_files:
            st.session_state.uploaded_files.append(uploaded_file.name)
    
    st.markdown("**Active Files:**")
    for file in st.session_state.uploaded_files:
        st.markdown(f"- 📄 {file}")
    
    st.markdown('</div>', unsafe_allow_html=True)

# ===== CHAT COLUMN =====
with chat_col:
    # Conditional rendering based on view_mode
    if st.session_state.view_mode == "email_view":
        # EMAIL LIST VIEW
        st.title(f"📧 Emails about {st.session_state.company_data['name']}")
        
        if st.button("← Back to Proposal", key="back_from_emails"):
            st.session_state.view_mode = "main"
            st.rerun()
        
        st.markdown("---")
        
        if st.session_state.company_emails:
            for i, email in enumerate(st.session_state.company_emails):
                with st.expander(f"📨 {email['subject']} - {email['sender']}", expanded=(i==0)):
                    st.markdown(f"**From:** {email['sender']}")
                    st.markdown(f"**Date:** {email['date']}")
                    st.markdown(f"**Subject:** {email['subject']}")
                    st.markdown("---")
                    st.markdown(email['body'])
        else:
            st.info("No emails found for this company.")
    
    elif st.session_state.view_mode == "chat_view":
        # TEAMS CHAT VIEW
        if st.session_state.company_chats:
            chat = st.session_state.company_chats[0]
            st.title(f"💬 {chat['title']}")
            
            if st.button("← Back to Proposal", key="back_from_chat"):
                st.session_state.view_mode = "main"
                st.rerun()
            
            st.markdown("---")
            st.markdown(f"**Participants:** {', '.join(chat['participants'])}")
            st.markdown("---")
            
            for msg in chat['messages']:
                st.markdown(f"**{msg['sender']}** - *{msg['timestamp']}*")
                st.markdown(msg['content'])
                st.markdown("")
        else:
            st.info("No group chats found for this company.")
    
    else:
        # MAIN PROPOSAL VIEW
        st.title("🤖 NexusCRM Sales Proposal Copilot")
    st.markdown("---")
    
    # Display chat messages
    for message in st.session_state.messages:
        avatar = "https://upload.wikimedia.org/wikipedia/en/a/aa/Microsoft_Copilot_Icon.svg" if message["role"] == "assistant" else "https://ui-avatars.com/api/?name=User&background=random"
        with st.chat_message(message["role"], avatar=avatar):
            st.markdown(message["content"])
            
            # Show draft editor if this message has it
            if message.get("show_editor"):
                with st.container():
                    st.markdown('<div class="draft-container">', unsafe_allow_html=True)
                    st.subheader("Edit Proposal Draft")
                    st.session_state.company_data["edited_full_draft"] = st.text_area(
                        "Edit the proposal content below. Use markdown headers for sections.", 
                        value=st.session_state.company_data["edited_full_draft"],
                        height=500,
                        key=f"full_draft_{message.get('id', 0)}"
                    )
                    
                    if st.button("✨ Generate PPT", key=f"confirm_{message.get('id', 0)}", use_container_width=True):
                        st.session_state.messages.append({
                            "role": "assistant",
                            "content": "Generating your PowerPoint presentation...",
                            "show_download": True,
                            "id": len(st.session_state.messages)
                        })
                        st.rerun()
                    
                    st.markdown('</div>', unsafe_allow_html=True)
            
            # Show download/preview button if this message has it
            if message.get("show_download"):
                with st.container():
                    st.markdown('<div class="draft-container">', unsafe_allow_html=True)
                    st.subheader("PPT Preview & Customization")

                    if (
                        st.session_state.last_sent_message_id == message.get('id', 0)
                        and st.session_state.last_sent_recipients
                    ):
                        sent_list = ", ".join(st.session_state.last_sent_recipients)
                        st.success(f"✅ Proposal sent to: {sent_list}")
                    
                    # PPT Viewer Simulation
                    tabs = st.tabs(["Slide 1", "Slide 2", "Slide 3", "Slide 4"])
                    
                    # Render mock slides in tabs
                    draft_text = st.session_state.company_data["edited_full_draft"]
                    # (Quick split for preview)
                    prev_sections = {"Executive Summary": "", "Solution": "", "Investment": ""}
                    curr = "Executive Summary"
                    for l in draft_text.split('\n'):
                        if "Executive Summary" in l: curr = "Executive Summary"
                        elif "Solution" in l: curr = "Solution"
                        elif "Investment" in l: curr = "Investment"
                        else: prev_sections[curr] += l + "\n"

                    theme = st.session_state.company_data["ppt_theme"]
                    bg_hex = '#%02x%02x%02x' % tuple(theme['bg_color'])
                    text_hex = '#%02x%02x%02x' % tuple(theme['title_color'])

                    def slide_style(content):
                        return f"""
                        <div style="background-color: {bg_hex}; border: 1px solid #ccc; padding: 20px; border-radius: 5px; min-height: 200px; color: {text_hex};">
                            <h3 style="color: {text_hex};">{content['title']}</h3>
                            <p style="color: #333; font-size: 14px;">{content['body'][:300]}...</p>
                        </div>
                        """

                    tabs[0].markdown(slide_style({"title": f"NexusCRM → {st.session_state.company_data['name']}", "body": "Strategic Proposal"}), unsafe_allow_html=True)
                    tabs[1].markdown(slide_style({"title": "Needs", "body": prev_sections["Executive Summary"]}), unsafe_allow_html=True)
                    tabs[2].markdown(slide_style({"title": "Solution", "body": prev_sections["Solution"]}), unsafe_allow_html=True)
                    tabs[3].markdown(slide_style({"title": "Investment", "body": prev_sections["Investment"]}), unsafe_allow_html=True)

                    st.markdown("---")
                    theme_suggestion = st.text_input("🎨 Suggest your theme changes", placeholder="e.g. Dark mode with gold accents", key=f"theme_input_{message.get('id', 0)}")
                    
                    cta_regen, cta_download, cta_send = st.columns([0.8, 1.1, 1.1])
                    if cta_regen.button("🔄 Regenerate Theme", key=f"regen_{message.get('id', 0)}", use_container_width=True):
                        if theme_suggestion:
                            with st.spinner("Applying theme changes..."):
                                new_theme = get_theme_update(theme_suggestion, st.session_state.company_data["ppt_theme"])
                                st.session_state.company_data["ppt_theme"] = new_theme
                                st.rerun()

                    pptx_path = None
                    try:
                        pptx_path = generate_pptx(st.session_state.company_data)
                    except Exception as e:
                        st.error(f"Error: {e}")

                    if pptx_path:
                        with open(pptx_path, "rb") as f:
                            cta_download.download_button(
                                "📥 Download Final (PPTX)",
                                f,
                                file_name=f"NexusCRM_Proposal_{st.session_state.company_data['name']}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                key=f"download_{message.get('id', 0)}",
                                use_container_width=True
                            )

                    if cta_send.button("📧 Send Proposal", key=f"send_proposal_preview_{message.get('id', 0)}", type="primary", use_container_width=True):
                        st.session_state.show_send_modal = True
                        st.session_state.send_modal_message_id = message.get('id', 0)
                        contacts = extract_contacts(
                            st.session_state.company_emails,
                            st.session_state.company_chats
                        )
                        st.session_state.selected_recipients = [c['email'] for c in contacts]
                        st.rerun()

                    render_send_modal(message.get('id', 0))
                    
                    st.markdown('</div>', unsafe_allow_html=True)

    # Chat input
    if prompt := st.chat_input("Ask Copilot... (Try: @SPG create proposal for Tesla)"):
        # Add user message
        st.session_state.messages.append({"role": "user", "content": prompt})
        
        # Check for @SPG trigger
        spg_match = re.search(r"@SPG.*for\s+(.+)", prompt, re.IGNORECASE)
        
        if spg_match:
            company_name = spg_match.group(1).strip()
            st.session_state.company_data["name"] = company_name
            
            # Generate contextual emails and chats
            from email_generator import generate_emails
            from teams_generator import generate_team_chat
            
            st.session_state.company_emails = generate_emails(company_name)
            st.session_state.company_chats = [generate_team_chat(company_name)]
            
            # Show thinking message
            with st.chat_message("assistant", avatar="https://upload.wikimedia.org/wikipedia/en/a/aa/Microsoft_Copilot_Icon.svg"):
                # Status 1: Web Search
                with st.spinner(f"🔍 Researcher: Searching Web, your emails and chats for details on {company_name}..."):
                    try:
                        research_results = research_company(
                            company_name,
                            st.session_state.company_emails,
                            st.session_state.company_chats
                        )
                    except Exception as e:
                        st.error(f"Error during research: {e}")
                        research_results = {
                            "executive_summary": "Unable to complete research.",
                            "solution": "Unable to complete research.",
                            "pricing": "Unable to complete research."
                        }
                
                # Status 2: Reading Emails & Teams
                with st.spinner("📧 Reading Emails & Teams Logs..."):
                    st.write(f"**Found {len(st.session_state.company_emails)} emails and {len(st.session_state.company_chats)} group chat**")
                
                # Status 3: Analyzing Knowledge
                with st.spinner("📂 Analyzing Knowledge Base..."):
                    st.write(f"**Analyzing {len(st.session_state.uploaded_files)} files...**")
                
                full_draft = f"## Executive Summary\n{research_results.get('executive_summary', '')}\n\n"
                full_draft += f"## Solution\n{research_results.get('solution', '')}\n\n"
                full_draft += f"## Investment\n{research_results.get('pricing', '')}"
                
                st.session_state.company_data["full_draft"] = full_draft
                st.session_state.company_data["edited_full_draft"] = full_draft
                
                st.session_state.messages.append({
                    "role": "assistant",
                    "content": f"I've gathered insights on **{company_name}** from web research, your email/Teams history, and the knowledge base. Here's a draft proposal - please review and edit:",
                    "show_editor": True,
                    "id": len(st.session_state.messages)
                })
            
            st.rerun()
        else:
            # Generic response
            st.session_state.messages.append({
                "role": "assistant",
                "content": "👋 I'm the **NexusCRM Sales Proposal Agent**. Tag me with `@SPG create proposal for [Company Name]` to get started!"
            })
            st.rerun()
